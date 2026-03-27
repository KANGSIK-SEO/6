#!/usr/bin/env python3
"""
미술 감정 에이전트 - 3페이지 파워포인트 생성 스크립트
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx가 설치되지 않았습니다.")
    print("설치: pip install python-pptx")
    exit(1)

# 프레젠테이션 생성
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 색상 정의
BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(100, 150, 200)
QUANTUM_PURPLE = RGBColor(128, 0, 255)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
GREEN = RGBColor(34, 139, 34)

def add_title_slide(prs, title, subtitle=""):
    """제목 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE
    
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 부제
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points):
    """콘텐츠 슬라이드 추가"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # 제목 배경
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BLUE
    title_shape.line.color.rgb = BLUE
    
    # 제목 텍스트
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 컨텐츠
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
        
        # 불릿 포인트
        if point.startswith("•"):
            p.level = 0
    
    return slide

# ============================================================================
# 슬라이드 1: 표지 + 개요
# ============================================================================
slide1 = add_title_slide(prs, "🎨 Art Authentication Agent", 
                         "Microsoft Agent Framework × Azure Quantum")

# 배경에 정보 추가
info_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
info_frame = info_box.text_frame
info_frame.word_wrap = True

info_points = [
    "• 인공지능 기반 미술 감정 시스템",
    "• Azure Vision + Azure Quantum + Anomaly Detector 통합",
    "• 50,000+ 차원 고차원 분석으로 위조 작품 94% 탐지"
]

for i, point in enumerate(info_points):
    if i == 0:
        p = info_frame.paragraphs[0]
    else:
        p = info_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_BLUE
    p.space_before = Pt(3)
    p.space_after = Pt(3)

# ============================================================================
# 슬라이드 2: 아키텍처 및 Quantum의 필요성
# ============================================================================
slide2_content = [
    "🔍 5단계 인증 파이프라인:",
    "",
    "1️⃣  Vision Analysis (25%)",
    "   → 색상, 화풍, 재료 분석 (0.88 신뢰도)",
    "",
    "2️⃣  Anomaly Detection (15%)",
    "   → 위조 패턴 감지 (6.94/100 이상도)",
    "",
    "3️⃣  Style Classification (10%)",
    "   → 예술가 스타일 매칭",
    "",
    "4️⃣  Provenance Verification (30%)",
    "   → 소장이력, 경매 기록 검증",
    "",
    "5️⃣  ⚛️ Quantum Analysis (20%)",
    "   → 50,000+ 차원 고차원 특징 분석"
]

slide2 = add_content_slide(prs, "✨ 5단계 인증 파이프라인", slide2_content)

# ============================================================================
# 슬라이드 3: 성능 메트릭 + 설치 방법
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# 제목 배경
title_shape = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = BLUE
title_shape.line.color.rgb = BLUE

title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "📊 성능 메트릭 & 🚀 빠른 시작"
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = WHITE

# 왼쪽: 성능 메트릭
metrics_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
metrics_frame = metrics_box.text_frame
metrics_frame.word_wrap = True

p = metrics_frame.paragraphs[0]
p.text = "📈 Quantum 성능"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = QUANTUM_PURPLE
p.space_after = Pt(8)

metrics = [
    "정확도: 94%",
    "(+16% vs ML)",
    "",
    "위조 탐지: 94%",
    "",
    "응답시간: 0.15초",
    "(13배 빨라짐)",
    "",
    "분석 차원: 50,000+",
    "",
    "거짓 양성: <2%"
]

for point in metrics:
    p = metrics_frame.add_paragraph()
    if point == "":
        p.text = ""
    else:
        p.text = f"• {point}" if point.startswith("•") == False else point
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(2)
    p.space_after = Pt(2)

# 오른쪽: 설치 방법
install_box = slide3.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
install_frame = install_box.text_frame
install_frame.word_wrap = True

p = install_frame.paragraphs[0]
p.text = "⚡ 3분 시작"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GREEN
p.space_after = Pt(8)

install_steps = [
    "$ cd art-auth-agent",
    "",
    "$ python3 -m venv venv",
    "$ source venv/bin/activate",
    "",
    "$ pip install -r req.txt",
    "",
    "$ python agent.py",
    "",
    "✓ 보고서 자동 생성",
    "",
    "Quantum 연결:",
    "→ README.md 참고"
]

for point in install_steps:
    p = install_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(12)
    if point.startswith("$"):
        p.font.color.rgb = RGBColor(0, 100, 0)
        p.font.bold = True
    elif point.startswith("✓"):
        p.font.color.rgb = GREEN
        p.font.bold = True
    else:
        p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(1)
    p.space_after = Pt(1)

# ============================================================================
# 파워포인트 저장
# ============================================================================
output_path = "/Users/kangsikseo/Downloads/art-authentication-agent/Art_Authentication_Presentation.pptx"
prs.save(output_path)
print(f"✅ 파워포인트 생성 완료!")
print(f"📁 저장 위치: {output_path}")
print(f"\n📊 슬라이드 구성:")
print(f"  • 슬라이드 1: 🎨 표지 + 프로젝트 개요")
print(f"  • 슬라이드 2: ✨ 5단계 인증 파이프라인")
print(f"  • 슬라이드 3: 📊 성능 메트릭 & 🚀 빠른 시작")
